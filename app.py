import os
import subprocess
from flask import Flask, render_template, request, redirect, url_for, flash, send_from_directory
from werkzeug.utils import secure_filename
from PIL import Image
from pdf2image import convert_from_path
import pytesseract

pytesseract.pytesseract.tesseract_cmd = r"your file path"

app = Flask(__name__)


# Create a folder for converted files if it doesn't exist.
UPLOAD_FOLDER = os.path.join(os.getcwd(), 'static', 'converted')
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# Mapping of conversion types to allowed file extensions
conversion_allowed = {
    "jpg_to_pdf": {"jpg", "jpeg", "png"},
    "pdf_to_jpg": {"pdf"},
    "ppt_to_pdf": {"ppt", "pptx"},
    "word_to_pdf": {"doc", "docx"},
    "excel_to_pdf": {"xls", "xlsx"},
    "pdf_ocr": {"pdf"},
    "pdf_to_word": {"pdf"},
    "pdf_to_ppt": {"pdf"}
}

def is_allowed_for_conversion(conversion_type, filename):
    if '.' not in filename:
        return False
    ext = filename.rsplit('.', 1)[1].lower()
    allowed_exts = conversion_allowed.get(conversion_type, set())
    return ext in allowed_exts

# -------------------------------
# Conversion Functions
# -------------------------------

def convert_jpg_to_pdf(file_stream, filename):
    try:
        image = Image.open(file_stream)
        if image.mode in ("RGBA", "P"):
            image = image.convert("RGB")
        out_name = filename.rsplit('.', 1)[0] + ".pdf"
        pdf_filename = secure_filename(out_name)
        pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], pdf_filename)
        image.save(pdf_path, "PDF", resolution=100.0)
        return [pdf_filename]
    except Exception as e:
        print("Error converting JPG to PDF:", e)
        return None

def convert_pdf_to_jpg(file_stream, filename):
    try:
        temp_pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(filename))
        file_stream.save(temp_pdf_path)
        poppler_path = r"your file path"
        images = convert_from_path(temp_pdf_path, poppler_path=poppler_path)
        output_files = []
        for i, image in enumerate(images):
            out_name = f"{filename.rsplit('.', 1)[0]}_page_{i+1}.jpg"
            out_filename = secure_filename(out_name)
            out_path = os.path.join(app.config['UPLOAD_FOLDER'], out_filename)
            image.save(out_path, 'JPEG')
            output_files.append(out_filename)
        os.remove(temp_pdf_path)
        return output_files
    except Exception as e:
        print("Error converting PDF to JPG:", e)
        return None

def ppt_to_pdf(file_stream, filename):
    try:
        temp_path = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(filename))
        file_stream.save(temp_path)
        out_dir = app.config['UPLOAD_FOLDER']
        libreoffice_path = r"your file path"
        command = [libreoffice_path, "--headless", "--convert-to", "pdf", "--outdir", out_dir, temp_path]
        subprocess.run(command, check=True)
        out_filename = secure_filename(filename.rsplit('.', 1)[0] + ".pdf")
        os.remove(temp_path)
        return [out_filename]
    except Exception as e:
        print("Error converting PPT to PDF:", e)
        return None

def word_to_pdf(file_stream, filename):
    try:
        temp_path = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(filename))
        file_stream.save(temp_path)
        out_dir = app.config['UPLOAD_FOLDER']
        libreoffice_path = r"your file path"
        command = [libreoffice_path, "--headless", "--convert-to", "pdf", "--outdir", out_dir, temp_path]
        subprocess.run(command, check=True)
        out_filename = secure_filename(filename.rsplit('.', 1)[0] + ".pdf")
        os.remove(temp_path)
        return [out_filename]
    except Exception as e:
        print("Error converting Word to PDF:", e)
        return None

def excel_to_pdf(file_stream, filename):
    try:
        temp_path = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(filename))
        file_stream.save(temp_path)
        out_dir = app.config['UPLOAD_FOLDER']
        libreoffice_path = r"your file path"
        command = [libreoffice_path, "--headless", "--convert-to", "pdf", "--outdir", out_dir, temp_path]
        subprocess.run(command, check=True)
        out_filename = secure_filename(filename.rsplit('.', 1)[0] + ".pdf")
        os.remove(temp_path)
        return [out_filename]
    except Exception as e:
        print("Error converting Excel to PDF:", e)
        return None

def pdf_ocr(file_stream, filename):
    try:
        temp_pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(filename))
        file_stream.save(temp_pdf_path)
        poppler_path = r"your file path"
        images = convert_from_path(temp_pdf_path, poppler_path=poppler_path)
        ocr_text = ""
        for i, image in enumerate(images):
            text = pytesseract.image_to_string(image)
            ocr_text += f"Page {i+1}:\n{text}\n\n"
        os.remove(temp_pdf_path)
        out_name = filename.rsplit('.', 1)[0] + "_ocr.txt"
        out_filename = secure_filename(out_name)
        out_path = os.path.join(app.config['UPLOAD_FOLDER'], out_filename)
        with open(out_path, "w", encoding="utf-8") as f:
            f.write(ocr_text)
        return [out_filename]
    except Exception as e:
        print("Error performing OCR on PDF:", e)
        return None

def pdf_to_word(file_stream, filename):
    try:
        temp_pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(filename))
        file_stream.save(temp_pdf_path)
        out_name = filename.rsplit('.', 1)[0] + ".docx"
        out_filename = secure_filename(out_name)
        out_path = os.path.join(app.config['UPLOAD_FOLDER'], out_filename)
        from pdf2docx import Converter
        cv = Converter(temp_pdf_path)
        cv.convert(out_path, start=0, end=None)
        cv.close()
        os.remove(temp_pdf_path)
        return [out_filename]
    except Exception as e:
        print("Error converting PDF to Word:", e)
        return None

def pdf_to_ppt(file_stream, filename):
    try:
        temp_pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(filename))
        file_stream.save(temp_pdf_path)
        poppler_path = r"your file path"
        images = convert_from_path(temp_pdf_path, poppler_path=poppler_path)
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        for image in images:
            slide = prs.slides.add_slide(blank_slide_layout)
            temp_img_path = os.path.join(app.config['UPLOAD_FOLDER'], "temp_slide.jpg")
            image.save(temp_img_path, 'JPEG')
            slide.shapes.add_picture(temp_img_path, Inches(0), Inches(0),
                                     width=prs.slide_width, height=prs.slide_height)
            os.remove(temp_img_path)
        out_name = filename.rsplit('.', 1)[0] + ".pptx"
        out_filename = secure_filename(out_name)
        out_path = os.path.join(app.config['UPLOAD_FOLDER'], out_filename)
        prs.save(out_path)
        os.remove(temp_pdf_path)
        return [out_filename]
    except Exception as e:
        print("Error converting PDF to PPT:", e)
        return None

# -------------------------------
# Flask Routes
# -------------------------------

@app.route('/')
def index():
    return render_template("index.html")

@app.route('/convert', methods=['POST'])
def convert():
    if 'file' not in request.files:
        flash("No file part")
        return redirect(url_for('index'))
    file = request.files['file']
    if file.filename == '':
        flash("No file selected")
        return redirect(url_for('index'))
    
    conversion_type = request.form.get("conversion_type")
    filename = secure_filename(file.filename)
    
    if not is_allowed_for_conversion(conversion_type, filename):
        flash("Invalid file type for this conversion.")
        return redirect(url_for('index'))
    
    flash("File uploaded successfully!")
    
    result = None
    if conversion_type == "jpg_to_pdf":
        result = convert_jpg_to_pdf(file, filename)
    elif conversion_type == "pdf_to_jpg":
        result = convert_pdf_to_jpg(file, filename)
    elif conversion_type == "ppt_to_pdf":
        result = ppt_to_pdf(file, filename)
    elif conversion_type == "word_to_pdf":
        result = word_to_pdf(file, filename)
    elif conversion_type == "excel_to_pdf":
        result = excel_to_pdf(file, filename)
    elif conversion_type == "pdf_ocr":
        result = pdf_ocr(file, filename)
    elif conversion_type == "pdf_to_word":
        result = pdf_to_word(file, filename)
    elif conversion_type == "pdf_to_ppt":
        result = pdf_to_ppt(file, filename)
    else:
        flash("Selected conversion type not implemented.")
        return redirect(url_for('index'))
    
    if result:
        download_files = []
        for fn in result:
            ext = fn.rsplit('.', 1)[-1].lower()
            file_info = {"filename": fn, "extension": ext}
            if ext == "txt":
                try:
                    with open(os.path.join(app.config['UPLOAD_FOLDER'], fn), "r", encoding="utf-8") as f:
                        file_info["preview_content"] = f.read()
                except Exception as e:
                    file_info["preview_content"] = "Could not load preview."
            download_files.append(file_info)
        msg = f"Conversion Successful: {conversion_type.replace('_', ' ').title()}"
        return render_template("result.html", download_files=download_files, message=msg)
    else:
        flash("Conversion failed or not implemented.")
        return redirect(url_for('index'))

@app.route('/preview/<filename>')
def preview_file(filename):
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename)

@app.route('/download/<filename>')
def download_file(filename):
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename, as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)
